"""
Corporate Template Analysis Tool
Extracts structural information from a PowerPoint template.
Outputs only METADATA - no proprietary content.

Usage: python analyze_template.py path/to/template.pptx
"""

import sys
import json
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

def analyze_template(template_path):
    """Complete template analysis - outputs only structural metadata."""
    
    prs = Presentation(template_path)
    
    analysis = {
        "file_info": analyze_file_info(template_path),
        "dimensions": analyze_dimensions(prs),
        "theme": analyze_theme(template_path),
        "slide_masters": analyze_masters(prs),
        "slide_layouts": analyze_layouts(prs),
        "slides": analyze_slides(prs),
        "media_inventory": analyze_media(template_path),
        "fonts_used": analyze_fonts(template_path),
        "colors_used": analyze_colors(template_path),
        "patterns_identified": identify_patterns(prs),
        "branding_elements": identify_branding(prs),
    }
    
    return analysis


def analyze_file_info(path):
    """Basic file information."""
    p = Path(path)
    return {
        "filename": p.name,
        "size_mb": round(p.stat().st_size / (1024 * 1024), 2),
    }


def analyze_dimensions(prs):
    """Slide canvas dimensions."""
    return {
        "width_inches": round(Emu(prs.slide_width).inches, 2),
        "height_inches": round(Emu(prs.slide_height).inches, 2),
        "aspect_ratio": f"{round(prs.slide_width / prs.slide_height, 2)}:1",
        "format": "Widescreen 16:9" if abs(prs.slide_width / prs.slide_height - 16/9) < 0.1 else 
                  "Standard 4:3" if abs(prs.slide_width / prs.slide_height - 4/3) < 0.1 else
                  "Custom",
    }


def analyze_theme(template_path):
    """Extract theme information from theme1.xml."""
    theme_info = {"colors": {}, "fonts": {}}
    
    with zipfile.ZipFile(template_path) as z:
        try:
            theme_xml = z.read("ppt/theme/theme1.xml")
            tree = etree.fromstring(theme_xml)
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            
            # Extract color scheme
            color_scheme = tree.find(".//a:clrScheme", ns)
            if color_scheme is not None:
                for color_elem in color_scheme:
                    color_name = etree.QName(color_elem).localname
                    for child in color_elem:
                        tag = etree.QName(child).localname
                        if tag == "srgbClr":
                            theme_info["colors"][color_name] = "#" + child.get("val")
                        elif tag == "sysClr":
                            theme_info["colors"][color_name] = child.get("lastClr", "system")
            
            # Extract font scheme
            font_scheme = tree.find(".//a:fontScheme", ns)
            if font_scheme is not None:
                major = font_scheme.find(".//a:majorFont/a:latin", ns)
                minor = font_scheme.find(".//a:minorFont/a:latin", ns)
                if major is not None:
                    theme_info["fonts"]["major"] = major.get("typeface")
                if minor is not None:
                    theme_info["fonts"]["minor"] = minor.get("typeface")
        except Exception as e:
            theme_info["error"] = str(e)
    
    return theme_info


def analyze_masters(prs):
    """Analyze slide masters."""
    masters = []
    for mi, master in enumerate(prs.slide_masters):
        master_info = {
            "index": mi,
            "layout_count": len(master.slide_layouts),
            "placeholder_count": len(master.placeholders),
            "shape_count": len(master.shapes),
        }
        masters.append(master_info)
    return masters


def analyze_layouts(prs):
    """Analyze available slide layouts."""
    layouts = []
    for master in prs.slide_masters:
        for li, layout in enumerate(master.slide_layouts):
            layout_info = {
                "index": li,
                "name": layout.name,
                "placeholder_count": len(layout.placeholders),
                "shape_count": len(layout.shapes),
                "placeholders": []
            }
            for ph in layout.placeholders:
                try:
                    ph_info = {
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type),
                        "name": ph.name,
                        "position_inches": {
                            "left": round(Emu(ph.left).inches, 2) if ph.left else 0,
                            "top": round(Emu(ph.top).inches, 2) if ph.top else 0,
                            "width": round(Emu(ph.width).inches, 2) if ph.width else 0,
                            "height": round(Emu(ph.height).inches, 2) if ph.height else 0,
                        }
                    }
                    layout_info["placeholders"].append(ph_info)
                except Exception:
                    pass
            layouts.append(layout_info)
    return layouts


def analyze_slides(prs):
    """Analyze each slide structurally - no content extracted."""
    slides = []
    for si, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": si + 1,
            "layout_used": slide.slide_layout.name if slide.slide_layout else None,
            "shape_count": len(slide.shapes),
            "shapes_by_type": {},
            "text_box_count": 0,
            "image_count": 0,
            "chart_count": 0,
            "table_count": 0,
            "group_count": 0,
            "freeform_count": 0,
            "placeholders_used": 0,
            "shape_details": []
        }
        
        for shape in slide.shapes:
            try:
                shape_type = str(shape.shape_type).split(".")[-1] if shape.shape_type else "UNKNOWN"
                slide_info["shapes_by_type"][shape_type] = slide_info["shapes_by_type"].get(shape_type, 0) + 1
                
                # Categorize
                if shape.has_text_frame:
                    slide_info["text_box_count"] += 1
                if shape.shape_type == 13:  # PICTURE
                    slide_info["image_count"] += 1
                if shape.has_chart:
                    slide_info["chart_count"] += 1
                if shape.has_table:
                    slide_info["table_count"] += 1
                if "Group" in shape.name:
                    slide_info["group_count"] += 1
                if "Freeform" in shape.name:
                    slide_info["freeform_count"] += 1
                if shape.is_placeholder:
                    slide_info["placeholders_used"] += 1
                
                # Position info (no content)
                shape_detail = {
                    "name": shape.name,
                    "type": shape_type,
                    "position": {
                        "left": round(Emu(shape.left).inches, 2) if shape.left else 0,
                        "top": round(Emu(shape.top).inches, 2) if shape.top else 0,
                        "width": round(Emu(shape.width).inches, 2) if shape.width else 0,
                        "height": round(Emu(shape.height).inches, 2) if shape.height else 0,
                    },
                    "is_placeholder": shape.is_placeholder,
                    "has_text": shape.has_text_frame,
                    "text_length": len(shape.text_frame.text) if shape.has_text_frame else 0,
                    # We extract LENGTH only, not actual text
                }
                slide_info["shape_details"].append(shape_detail)
            except Exception:
                pass
        
        slides.append(slide_info)
    return slides


def analyze_media(template_path):
    """Inventory of media files - no actual media extracted."""
    media = {"images": 0, "videos": 0, "audio": 0, "image_formats": set()}
    
    with zipfile.ZipFile(template_path) as z:
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                ext = name.split(".")[-1].lower()
                if ext in ["png", "jpg", "jpeg", "gif", "bmp", "svg", "wmf", "emf"]:
                    media["images"] += 1
                    media["image_formats"].add(ext)
                elif ext in ["mp4", "avi", "wmv", "mov"]:
                    media["videos"] += 1
                elif ext in ["mp3", "wav", "wma"]:
                    media["audio"] += 1
    
    media["image_formats"] = list(media["image_formats"])
    return media


def analyze_fonts(template_path):
    """Inventory of fonts used in the template."""
    fonts_used = set()
    
    with zipfile.ZipFile(template_path) as z:
        for name in z.namelist():
            if name.endswith(".xml") and ("slide" in name or "theme" in name or "Master" in name):
                try:
                    content = z.read(name).decode("utf-8", errors="ignore")
                    # Find all typeface attributes
                    import re
                    matches = re.findall(r'typeface="([^"]+)"', content)
                    fonts_used.update(matches)
                except Exception:
                    pass
    
    return sorted(list(fonts_used))


def analyze_colors(template_path):
    """Inventory of hex colors used in the template."""
    colors_used = {}
    
    with zipfile.ZipFile(template_path) as z:
        for name in z.namelist():
            if name.endswith(".xml") and "slide" in name:
                try:
                    content = z.read(name).decode("utf-8", errors="ignore")
                    import re
                    # Find all srgbClr values
                    matches = re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', content)
                    for hex_color in matches:
                        hex_upper = hex_color.upper()
                        colors_used[hex_upper] = colors_used.get(hex_upper, 0) + 1
                except Exception:
                    pass
    
    # Return top 20 most-used colors
    sorted_colors = sorted(colors_used.items(), key=lambda x: -x[1])[:20]
    return [{"hex": "#" + c, "occurrence_count": n} for c, n in sorted_colors]


def identify_patterns(prs):
    """Identify visual patterns across slides without reading content."""
    patterns = []
    
    for si, slide in enumerate(prs.slides):
        pattern = {
            "slide_number": si + 1,
            "pattern_id": None,
            "description": "",
            "structural_signature": ""
        }
        
        # Count shape characteristics
        groups = sum(1 for s in slide.shapes if "Group" in s.name)
        text_boxes = sum(1 for s in slide.shapes if s.has_text_frame)
        images = sum(1 for s in slide.shapes if s.shape_type == 13)
        freeforms = sum(1 for s in slide.shapes if "Freeform" in s.name)
        
        # Build structural signature (no content)
        pattern["structural_signature"] = f"groups={groups}, texts={text_boxes}, images={images}, freeforms={freeforms}"
        
        # Identify pattern type by structure
        if si == 0:
            pattern["pattern_id"] = "title_slide"
            pattern["description"] = "Title slide pattern"
        elif groups == 3 and text_boxes >= 6:
            pattern["pattern_id"] = "three_column_cards"
            pattern["description"] = "3-column card layout"
        elif groups == 4 and text_boxes >= 8:
            pattern["pattern_id"] = "four_quadrant"
            pattern["description"] = "4-quadrant layout"
        elif images >= 1 and groups >= 2:
            pattern["pattern_id"] = "image_with_topics"
            pattern["description"] = "Image with content topics"
        elif text_boxes >= 4 and groups == 0:
            pattern["pattern_id"] = "title_with_content"
            pattern["description"] = "Title with content layout"
        else:
            pattern["pattern_id"] = "custom"
            pattern["description"] = "Custom layout"
        
        patterns.append(pattern)
    
    return patterns


def identify_branding(prs):
    """Identify recurring branding elements across slides."""
    branding = {
        "recurring_shapes": [],
        "background_pattern": "unknown",
        "header_footer_present": False,
        "logo_likely_present": False
    }
    
    # Look for shapes that appear on multiple slides at same position
    shape_positions = {}
    
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.left and shape.top:
                    key = (shape.name.split()[0], 
                           round(Emu(shape.left).inches, 1), 
                           round(Emu(shape.top).inches, 1))
                    shape_positions[key] = shape_positions.get(key, 0) + 1
            except Exception:
                pass
    
    # Recurring shapes (appear on 50%+ of slides)
    total_slides = len(prs.slides)
    for key, count in shape_positions.items():
        if count >= total_slides * 0.5:
            branding["recurring_shapes"].append({
                "shape_type": key[0],
                "position_inches": {"left": key[1], "top": key[2]},
                "appears_on_slides_pct": round(100 * count / total_slides, 1)
            })
    
    # Detect background pattern
    backgrounds = sum(1 for slide in prs.slides 
                      for shape in slide.shapes 
                      if "Freeform" in shape.name and shape.width and 
                      Emu(shape.width).inches > 15)
    if backgrounds >= total_slides * 0.8:
        branding["background_pattern"] = "consistent_background_shape"
    
    return branding


def main():
    if len(sys.argv) < 2:
        print("Usage: python analyze_template.py <template.pptx>")
        sys.exit(1)
    
    template_path = sys.argv[1]
    
    if not Path(template_path).exists():
        print(f"Error: File not found: {template_path}")
        sys.exit(1)
    
    print(f"Analyzing: {template_path}")
    print("=" * 70)
    
    analysis = analyze_template(template_path)
    
    # Output as JSON
    output_path = Path(template_path).stem + "_analysis.json"
    with open(output_path, "w") as f:
        json.dump(analysis, f, indent=2, default=str)
    
    print(f"\nAnalysis saved to: {output_path}")
    print("\n--- SUMMARY ---")
    print(f"Dimensions: {analysis['dimensions']['width_inches']}\" × {analysis['dimensions']['height_inches']}\"")
    print(f"Format: {analysis['dimensions']['format']}")
    print(f"Slides: {len(analysis['slides'])}")
    print(f"Layouts available: {len(analysis['slide_layouts'])}")
    print(f"Theme colors: {len(analysis['theme']['colors'])}")
    print(f"Major font: {analysis['theme']['fonts'].get('major', 'N/A')}")
    print(f"Minor font: {analysis['theme']['fonts'].get('minor', 'N/A')}")
    print(f"Fonts used: {', '.join(analysis['fonts_used'][:5])}{'...' if len(analysis['fonts_used']) > 5 else ''}")
    print(f"Top colors: {', '.join([c['hex'] for c in analysis['colors_used'][:5]])}")
    print(f"Images embedded: {analysis['media_inventory']['images']}")
    print(f"\nPatterns identified:")
    for p in analysis['patterns_identified']:
        print(f"  Slide {p['slide_number']}: {p['pattern_id']} - {p['description']}")
    print(f"\nBranding:")
    print(f"  Recurring shapes: {len(analysis['branding_elements']['recurring_shapes'])}")
    print(f"  Background pattern: {analysis['branding_elements']['background_pattern']}")


if __name__ == "__main__":
    main()




second 

def analyze_backgrounds(template_path):
    """Identify background shapes across slides."""
    prs = Presentation(template_path)
    canvas_w = Emu(prs.slide_width).inches
    canvas_h = Emu(prs.slide_height).inches
    
    background_analysis = {
        "slide_dimensions": {"width": canvas_w, "height": canvas_h},
        "per_slide_backgrounds": [],
        "recurring_background_shapes": [],
        "background_complexity": "unknown"
    }
    
    # Track shape positions across slides
    shape_signatures = {}
    
    for si, slide in enumerate(prs.slides):
        slide_bg = {
            "slide_number": si + 1,
            "background_shapes": []
        }
        
        # Process shapes in z-order (first added = bottom)
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                if not (shape.left and shape.top and shape.width and shape.height):
                    continue
                
                left = Emu(shape.left).inches
                top = Emu(shape.top).inches
                w = Emu(shape.width).inches
                h = Emu(shape.height).inches
                
                # Heuristic: is this a background shape?
                # 1. Large coverage (>50% of canvas)
                # 2. No text content
                # 3. Early in z-order (first few shapes)
                # 4. Doesn't overlap with content area
                
                area_coverage = (w * h) / (canvas_w * canvas_h)
                is_large = area_coverage > 0.30
                is_decorative = not shape.has_text_frame or (shape.has_text_frame and not shape.text_frame.text)
                is_back_of_stack = shape_idx < 8  # First 8 shapes are usually background
                is_freeform_or_autoshape = "Freeform" in shape.name or "AutoShape" in shape.name or "Group" in shape.name
                
                if (is_large or is_freeform_or_autoshape) and is_decorative and is_back_of_stack:
                    bg_shape = {
                        "name": shape.name,
                        "type": str(shape.shape_type),
                        "position": {"left": round(left, 2), "top": round(top, 2)},
                        "size": {"width": round(w, 2), "height": round(h, 2)},
                        "coverage_pct": round(area_coverage * 100, 1),
                        "z_order": shape_idx,
                        "extends_beyond_slide": (left < 0 or top < 0 or 
                                                left + w > canvas_w or 
                                                top + h > canvas_h)
                    }
                    slide_bg["background_shapes"].append(bg_shape)
                    
                    # Track for recurring detection
                    sig = (shape.name.split()[0], round(left, 1), round(top, 1), 
                           round(w, 1), round(h, 1))
                    shape_signatures[sig] = shape_signatures.get(sig, 0) + 1
            except Exception:
                pass
        
        background_analysis["per_slide_backgrounds"].append(slide_bg)
    
    # Find recurring background shapes
    total_slides = len(prs.slides)
    for sig, count in shape_signatures.items():
        if count >= total_slides * 0.5:
            background_analysis["recurring_background_shapes"].append({
                "shape_type": sig[0],
                "position": {"left": sig[1], "top": sig[2]},
                "size": {"width": sig[3], "height": sig[4]},
                "appears_on_pct": round(100 * count / total_slides, 1)
            })
    
    # Classify complexity
    avg_bg_shapes = sum(len(s["background_shapes"]) for s in background_analysis["per_slide_backgrounds"]) / max(total_slides, 1)
    if avg_bg_shapes < 1:
        background_analysis["background_complexity"] = "simple (solid color or basic)"
    elif avg_bg_shapes < 3:
        background_analysis["background_complexity"] = "moderate (gradient or simple shapes)"
    elif avg_bg_shapes < 6:
        background_analysis["background_complexity"] = "complex (multiple decorative shapes)"
    else:
        background_analysis["background_complexity"] = "highly complex (many layered shapes)"
    
    return background_analysis