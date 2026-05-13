"""
Smoke test for the Capgemini layout router.

Generates one test slide per known layout using sample content. Open the
resulting test_output.pptx in PowerPoint and visually verify every layout
renders correctly with text fitting cleanly inside each placeholder.

Usage:
    python validate_template.py path/to/capgemini_template.pptx
"""

import sys
from pptx import Presentation

# Adjust these imports to match your project layout
from layout_router import LAYOUT_MAP, get_slot_map
from render_slide_content import render_slide_content


SAMPLE_CONTENT = {
    "title": "The 2024 Cloud Market Landscape",
    "chapter_label": "Market Overview",
    "subhead": "Three vendors capture 66% of a $79.8B market growing 21% YoY",
    "bullets": [
        "**AWS** holds 31% share but shows a slight downward trend",
        "**Microsoft Azure** doubled to 25% share over seven years",
        "**Google Cloud** is fastest-growing at 11%, recently profitable",
        "**Big Three** account for 66% of $79.8B Q1 2024 spending",
        "**Inflection**: GCP joins AWS and Azure in profitability",
    ],
    "bullets_right": [
        "**Hybrid adoption** accelerating across regulated industries",
        "**AI workloads** driving 35% of new cloud spending in 2024",
        "**Data sovereignty** rules reshaping vendor selection in EU/APAC",
    ],
    "bullets_third": [
        "**Cost optimization** is the #1 priority for 2025",
        "**FinOps** team adoption up 60% YoY",
        "**Reserved capacity** purchases shifting to 1-year commits",
    ],
    "author_line": "Talk2Docs Team | November 2025",
    "speaker_notes": (
        "This slide sets up the market context. Emphasize that while AWS is "
        "still the leader, the gap is narrowing. Azure's doubling over seven "
        "years is the key trend to remember."
    ),
}


class MockSpec:
    """Mimics a SlideSpec for testing without depending on pydantic."""
    def __init__(self, **kwargs):
        for k, v in kwargs.items():
            setattr(self, k, v)


def main(template_path: str, output_path: str = "test_output.pptx"):
    prs = Presentation(template_path)

    # Build a quick lookup: layout_name → slide_layout object
    layouts_by_name = {layout.name: layout for layout in prs.slide_layouts}

    # We want to add slides AFTER any existing slides in the template.
    # Some templates have demo slides; we leave them and append.
    existing_slide_count = len(prs.slides)
    print(f"Template has {existing_slide_count} existing slides, "
          f"{len(layouts_by_name)} layouts.")

    tested = []
    skipped = []

    for layout_name, slot_map in LAYOUT_MAP.items():
        if layout_name not in layouts_by_name:
            skipped.append((layout_name, "layout not in template"))
            continue

        layout = layouts_by_name[layout_name]
        slide = prs.slides.add_slide(layout)

        spec = MockSpec(
            title=f"{layout_name}: {SAMPLE_CONTENT['title']}",
            chapter_label=SAMPLE_CONTENT["chapter_label"],
            subhead=SAMPLE_CONTENT["subhead"],
            bullets=SAMPLE_CONTENT["bullets"],
            bullets_right=SAMPLE_CONTENT["bullets_right"],
            bullets_third=SAMPLE_CONTENT["bullets_third"],
            author_line=SAMPLE_CONTENT["author_line"],
            speaker_notes=SAMPLE_CONTENT["speaker_notes"],
        )

        try:
            render_slide_content(slide, spec, layout_name)
            tested.append(layout_name)
        except Exception as e:
            skipped.append((layout_name, f"render error: {e}"))

    prs.save(output_path)

    print(f"\nWrote {output_path}")
    print(f"Successfully rendered {len(tested)} layouts:")
    for name in tested:
        print(f"  ✓ {name}")
    if skipped:
        print(f"\nSkipped {len(skipped)} layouts:")
        for name, reason in skipped:
            print(f"  ✗ {name}: {reason}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python validate_template.py <template.pptx> [output.pptx]")
        sys.exit(1)
    template = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else "test_output.pptx"
    main(template, output)
