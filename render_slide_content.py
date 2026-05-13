# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""
Presentation orchestrator for the Capgemini Talk2Docs brand-aligned
presentation agent.

This module owns the end-to-end render pipeline:
  1. Resolve the user's PPTX template (local path or GCS)
  2. Strip any existing demo slides
  3. Render the cover slide
  4. Render every body slide, using the layout_router to place content into
     the correct placeholders for the specific Capgemini layout chosen
  5. Render the closing slide
  6. Save to disk and return the path

The key change vs the previous version is that render_slide_content now
delegates placeholder routing to layout_router.get_placeholder_for_role(),
so content lands in the right slot for each of the template's 59 layouts
(including the tricky ones where the TITLE placeholder is decorative).
"""

import asyncio
import json
import os
import tempfile
import uuid

from google.adk.tools.tool_context import ToolContext
from google.genai import types
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt

from ..shared_libraries.config import (
    DEFAULT_TEMPLATE_URI,
    GCS_BUCKET_NAME,
    get_logger,
)
from ..shared_libraries.models import CoverSpec, DeckSpec, SlideSpec
from ..shared_libraries.utils import _insert_image
from .artifact_utils import get_gcs_file_as_local_path, save_presentation
from .layout_router import (
    LAYOUT_MAP,
    PREFERRED_LAYOUTS_FOR_SYNTHESIZER,
    SlotMap,
    get_placeholder_for_role,
    get_slot_map,
    is_layout_known,
)
from .pptx_editor import _insert_visual_into_slide
from .visual_generator import generate_visual


# =============================================================================
# Layout selection
# =============================================================================

def get_smart_layout(prs: Presentation, requested_name: str):
    """
    Map a requested layout name to the best matching layout in the template.

    Priority:
      1. Exact match against a known Capgemini layout name (preferred).
      2. Keyword-based mapping to the canonical Capgemini layouts.
      3. Conceptual fallbacks (title/content/closing).
      4. Last-resort: the first layout in the template.
    """
    log = get_logger("layout_mapper")
    if not requested_name:
        requested_name = "Content 1 Chapterbox"
    original_request = requested_name
    requested_name_lower = requested_name.lower()

    layouts = prs.slide_layouts

    # 1. Exact match (preferred path — synthesizer should produce canonical names)
    for layout in layouts:
        if layout.name == requested_name:
            return layout
    for layout in layouts:
        if layout.name.lower() == requested_name_lower:
            return layout

    # 2. Keyword-based mapping to Capgemini layouts.
    # Order matters — more specific concepts first.
    keyword_map = [
        # (keyword in request, preferred Capgemini layout name)
        ("cover", "Title Slide 1"),
        ("title slide", "Title Slide 1"),
        ("opening", "Title Slide 1"),

        ("section header", "Title Slide Message left"),
        ("divider", "Title Slide Message left"),
        ("transition", "Title Slide Message left"),

        ("comparison", "Comparison Chapterbox"),
        ("compare", "Comparison Chapterbox"),
        ("side by side", "Comparison Chapterbox"),
        ("two content", "Comparison Chapterbox"),

        ("three", "Content 3 Boxes Chapterbox"),
        ("3 boxes", "Content 3 Boxes Chapterbox"),

        ("agenda", "Agenda 3"),
        ("toc", "Agenda 3"),
        ("roadmap", "Agenda 3"),

        ("closing", "Conclusion 1"),
        ("thank you", "End Slide 1"),
        ("end slide", "End Slide 1"),
        ("contact", "Conclusion 1"),

        ("title only", "Title Only Chapterbox"),
        ("image", "Content 2 Chapterbox"),
        ("picture", "Content 2 Chapterbox"),
        ("photo", "Content 2 Chapterbox"),

        # Generic content fallbacks
        ("title and content", "Content 1 Chapterbox"),
        ("title subtitle", "Title Subtitle Chapterbox"),
        ("content", "Content 1 Chapterbox"),
        ("bullet", "Content 1 Chapterbox"),
    ]

    for keyword, target_layout_name in keyword_map:
        if keyword in requested_name_lower:
            for layout in layouts:
                if layout.name == target_layout_name:
                    log.info(
                        f"Mapped '{original_request}' -> '{target_layout_name}'"
                    )
                    return layout

    # 3. Last resort: pick the first Content 1 Chapterbox if present, else first layout
    for layout in layouts:
        if layout.name == "Content 1 Chapterbox":
            log.info(
                f"No match for '{original_request}', falling back to Content 1 Chapterbox"
            )
            return layout

    log.warning(
        f"No matching layout for '{original_request}', using first layout: {layouts[0].name}"
    )
    return layouts[0]


# =============================================================================
# Render helpers — used by render_slide_content
# =============================================================================

def _rm_md(t: str) -> str:
    """Strip basic markdown markers."""
    if not t:
        return ""
    return t.replace("**", "")


def _truncate(text: str, max_chars: int, suffix: str = "…") -> str:
    """Hard-truncate text at max_chars, preferring a word boundary."""
    if not text or len(text) <= max_chars:
        return text
    cut = text[: max_chars - len(suffix)]
    last_space = cut.rfind(" ")
    if last_space > max_chars * 0.6:
        cut = cut[:last_space]
    return cut + suffix


def _configure_text_frame(tf, autosize: bool = True):
    """Apply safe defaults so text fits its placeholder."""
    tf.word_wrap = True
    if autosize:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP


def _write_plain_text(placeholder, text: str, max_chars: int = None):
    """Write a single string into a placeholder safely."""
    if placeholder is None or not text:
        return
    if max_chars:
        text = _truncate(text, max_chars)
    placeholder.text = _rm_md(text)
    _configure_text_frame(placeholder.text_frame)


def _apply_bullets(text_frame, bullets, max_bullets: int = 6,
                   bullet_max_chars: int = 120):
    """Write a list of bullet strings, parsing **bold** segments."""
    text_frame.clear()
    _configure_text_frame(text_frame)

    bullets = list(bullets)[:max_bullets]

    for i, bullet_text in enumerate(bullets):
        is_sub = (
            bullet_text.startswith("  ")
            or bullet_text.startswith("\t")
            or bullet_text.startswith("- ")
        )

        clean_text = bullet_text.strip(" \t-•*")
        clean_text = _truncate(clean_text, bullet_max_chars)

        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.level = 1 if is_sub else 0

        parts = clean_text.split("**")
        for j, part in enumerate(parts):
            if not part:
                continue
            r = p.add_run()
            r.text = part
            if j % 2 != 0:
                r.font.bold = True


def _remove_placeholder(placeholder):
    """Remove an empty placeholder so master prompt text doesn't leak through."""
    if placeholder is None:
        return
    try:
        sp = placeholder._element
        sp.getparent().remove(sp)
    except Exception:
        pass


def _insert_picture_into_placeholder(placeholder, image_source):
    """Insert an image into a picture-typed placeholder."""
    if placeholder is None or not image_source:
        return False
    try:
        placeholder.insert_picture(image_source)
        return True
    except Exception:
        return False


# =============================================================================
# Main render function — Capgemini-template-aware
# =============================================================================

def render_slide_content(slide, spec_obj, layout_name: str, prs=None,
                         log=None, is_cover: bool = False):
    """
    Populate a slide based on its layout and the spec's content.

    Uses layout_router to find the correct placeholder for each semantic role
    (chapter_label, title, subhead, body, body_right, body_third, picture,
    cover_subtitle, cover_author), rather than guessing from placeholder types.

    Args:
        slide: a python-pptx Slide object created from a Capgemini layout.
        spec_obj: a SlideSpec, CoverSpec, or DeckCover/DeckClosing pydantic
                  model with fields: title, subhead, bullets, etc.
        layout_name: the layout name string used to create this slide.
                     Must match a key in layout_router.LAYOUT_MAP.
        prs: optional Presentation object (for floating images when no
             picture placeholder is available).
        log: optional logger.
        is_cover: kept for backwards-compatibility; when True, also infers
                  cover behaviour from the layout's slot_map.is_cover flag.
    """
    slot_map: SlotMap = get_slot_map(layout_name)

    # Treat as cover if either the caller said so or the layout says so.
    treat_as_cover = is_cover or slot_map.is_cover

    # ---------------------------------------------------------------
    # Short-circuit: layouts with no text content (End Slide, Blank)
    # ---------------------------------------------------------------
    if slot_map.is_end_slide or slot_map.is_blank:
        if log:
            log.info(f"Layout '{layout_name}' has no text slots — skipping.")
        return

    # ---------------------------------------------------------------
    # 1. CHAPTER LABEL (Chapterbox layouts only)
    # ---------------------------------------------------------------
    chapter_ph = get_placeholder_for_role(slide, layout_name, "chapter_label")
    chapter_text = getattr(spec_obj, "chapter_label", None)
    if chapter_ph is not None:
        if chapter_text:
            _write_plain_text(chapter_ph, chapter_text, max_chars=50)
        else:
            _remove_placeholder(chapter_ph)

    # ---------------------------------------------------------------
    # 2. TITLE
    # ---------------------------------------------------------------
    title_ph = get_placeholder_for_role(slide, layout_name, "title")
    title_text = getattr(spec_obj, "title", None)
    if title_ph is not None and title_text:
        _write_plain_text(title_ph, title_text,
                          max_chars=slot_map.title_max_chars)
        if treat_as_cover and title_ph.text_frame.paragraphs:
            title_ph.text_frame.paragraphs[0].alignment = (
                PP_PARAGRAPH_ALIGNMENT.CENTER
            )

    # ---------------------------------------------------------------
    # 3. SUBHEAD (blue 20pt accent line on workhorse layouts)
    # ---------------------------------------------------------------
    subhead_ph = get_placeholder_for_role(slide, layout_name, "subhead")
    subhead_text = getattr(spec_obj, "subhead", None)
    if subhead_ph is not None:
        if subhead_text:
            _write_plain_text(subhead_ph, subhead_text,
                              max_chars=slot_map.subhead_max_chars)
        else:
            _remove_placeholder(subhead_ph)

    # ---------------------------------------------------------------
    # 4. COVER-ONLY SLOTS (subtitle + author/date)
    # ---------------------------------------------------------------
    if treat_as_cover:
        cover_sub_ph = get_placeholder_for_role(
            slide, layout_name, "cover_subtitle"
        )
        if cover_sub_ph is not None:
            sub_text = (
                getattr(spec_obj, "subhead", None)
                or getattr(spec_obj, "cover_subtitle", None)
            )
            if sub_text:
                _write_plain_text(cover_sub_ph, sub_text, max_chars=60)
            else:
                _remove_placeholder(cover_sub_ph)

        cover_author_ph = get_placeholder_for_role(
            slide, layout_name, "cover_author"
        )
        if cover_author_ph is not None:
            author_text = getattr(spec_obj, "author_line", None)
            if author_text:
                _write_plain_text(cover_author_ph, author_text, max_chars=60)
            else:
                _remove_placeholder(cover_author_ph)

        # Optional cover background image
        image_source = (
            getattr(spec_obj, "image_data", None)
            or getattr(spec_obj, "image_file_path", None)
        )
        if image_source:
            picture_ph = get_placeholder_for_role(slide, layout_name, "picture")
            _insert_picture_into_placeholder(picture_ph, image_source)

        # Covers have no body content; stop here.
        return

    # ---------------------------------------------------------------
    # 5. BODY BULLETS (single, two-column, three-column)
    # ---------------------------------------------------------------
    has_bullets = (
        hasattr(spec_obj, "bullets") and bool(getattr(spec_obj, "bullets", None))
    )

    if has_bullets:
        body_ph = get_placeholder_for_role(slide, layout_name, "body")
        if body_ph is not None:
            _apply_bullets(
                body_ph.text_frame,
                spec_obj.bullets,
                max_bullets=slot_map.body_max_bullets,
                bullet_max_chars=slot_map.body_bullet_max_chars,
            )

    bullets_right = getattr(spec_obj, "bullets_right", None)
    if bullets_right:
        body_right_ph = get_placeholder_for_role(
            slide, layout_name, "body_right"
        )
        if body_right_ph is not None:
            _apply_bullets(
                body_right_ph.text_frame,
                bullets_right,
                max_bullets=slot_map.body_max_bullets,
                bullet_max_chars=slot_map.body_bullet_max_chars,
            )

    bullets_third = getattr(spec_obj, "bullets_third", None)
    if bullets_third:
        body_third_ph = get_placeholder_for_role(
            slide, layout_name, "body_third"
        )
        if body_third_ph is not None:
            _apply_bullets(
                body_third_ph.text_frame,
                bullets_third,
                max_bullets=slot_map.body_max_bullets,
                bullet_max_chars=slot_map.body_bullet_max_chars,
            )

    # ---------------------------------------------------------------
    # 6. IMAGE (if the layout supports one)
    # ---------------------------------------------------------------
    image_source = (
        getattr(spec_obj, "image_data", None)
        or getattr(spec_obj, "image_file_path", None)
    )
    if image_source:
        picture_ph = get_placeholder_for_role(slide, layout_name, "picture")
        if picture_ph is not None:
            inserted = _insert_picture_into_placeholder(
                picture_ph, image_source
            )
            if not inserted and log:
                log.warning(f"Failed to insert picture into '{layout_name}'")
        elif prs is not None:
            # Fallback: no picture placeholder on this layout but we have an
            # image to render — float it in a safe lower-right region.
            try:
                box_hint = (
                    int(prs.slide_width * 0.55),
                    int(prs.slide_height * 0.30),
                    int(prs.slide_width * 0.40),
                    int(prs.slide_height * 0.55),
                )
                _insert_image(prs, slide, image_source, box_hint=box_hint)
            except Exception as e:
                if log:
                    log.warning(f"Failed to float image: {e}")


# =============================================================================
# Visual layout coercion — used by generate_and_render_deck
# =============================================================================

# Layouts that explicitly include a picture placeholder. The synthesizer's
# layout choice will be overridden to one of these when it asks for a visual.
LAYOUTS_WITH_PICTURE = [
    "Content 2 Chapterbox",       # Content + sidebar (preferred for images)
    "Content 2",
    "Divider Photo 1",
    "Divider Photo 2",
    "Title Slide 2",
    "Title Slide 3",
    "Conclusion 1 Photo",
    "Conclusion 2 Photo",
]

DEFAULT_IMAGE_LAYOUT = "Content 2 Chapterbox"


# =============================================================================
# render_deck_from_spec — top-level renderer
# =============================================================================

async def render_deck_from_spec(
    spec_dict: dict,
    out_pptx: str,
    tool_context: ToolContext,
    template_pptx: str | None = None,
) -> str:
    """
    Render a presentation from a spec, using a template if provided.
    Saves to a temp file and returns the path.
    """
    log = get_logger("render_deck_from_spec")
    try:
        # 1. TEMPLATE SELECTION
        if template_pptx and os.path.exists(template_pptx):
            log.info(
                f"Using user template '{template_pptx}' as the foundation."
            )
            working_template = template_pptx
        else:
            log.error("No valid user template provided. Aborting.")
            return "Error: No valid template provided."

        prs = Presentation(working_template)

        # ---------------------------------------------------------------
        # 2. COVER
        # ---------------------------------------------------------------
        cover_data = spec_dict.get(
            "cover", {"title": "Strategic Research & Analysis"}
        )
        cover_spec = CoverSpec(**cover_data)

        # Strip all but the first existing slide (keep first as our cover host).
        if len(prs.slides) > 0:
            for i in range(len(prs.slides) - 1, 0, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]

        cover_layout_name = getattr(cover_spec, "layout_name", None) or "Title Slide 1"

        if len(prs.slides) > 0:
            log.info("Template has an existing slide. Using it as the Cover Page.")
            cover_slide = prs.slides[0]
            # Best-effort: figure out the existing slide's layout name
            try:
                existing_layout_name = cover_slide.slide_layout.name
                if is_layout_known(existing_layout_name):
                    cover_layout_name = existing_layout_name
            except Exception:
                pass
            try:
                render_slide_content(
                    cover_slide, cover_spec, cover_layout_name,
                    prs=prs, log=log, is_cover=True,
                )
            except Exception as e:
                log.warning(f"Could not render cover on existing slide: {e}")
        else:
            log.info("Template is empty. Generating a new Cover Page.")
            try:
                cover_slide = prs.slides.add_slide(
                    get_smart_layout(prs, cover_layout_name)
                )
                render_slide_content(
                    cover_slide, cover_spec, cover_layout_name,
                    prs=prs, log=log, is_cover=True,
                )
            except Exception as e:
                log.warning(f"Could not generate/render cover slide: {e}")

        # ---------------------------------------------------------------
        # 3. BODY SLIDES
        # ---------------------------------------------------------------
        for s_data in spec_dict.get("slides", []):
            if "title" not in s_data or not s_data["title"]:
                s_data["title"] = "Slide Content"

            try:
                s_spec = SlideSpec(**s_data)
                layout = get_smart_layout(prs, s_spec.layout_name)
                slide = prs.slides.add_slide(layout)
                # Use the actual layout name (post-mapping) so router picks
                # the right slots.
                actual_layout_name = layout.name
                render_slide_content(
                    slide, s_spec, actual_layout_name,
                    prs=prs, log=log,
                )
            except Exception as e:
                log.error(
                    f"Failed to render slide '{s_data.get('title')}': {e}"
                )
                continue

            # Speaker notes + citations on the notes slide.
            # (render_slide_content does not handle notes — kept here so
            # citation formatting stays unchanged from the previous version.)
            if getattr(s_spec, "speaker_notes", None) or getattr(s_spec, "citations", None):
                try:
                    notes = slide.notes_slide.notes_text_frame
                    text = s_spec.speaker_notes or ""
                    if s_spec.citations:
                        if text:
                            text += "\n\n---\nCitations:\n"
                        else:
                            text += "Citations:\n"
                        for citation in s_spec.citations:
                            text += f"- {citation}\n"
                    notes.text = text
                except Exception:
                    pass

        # ---------------------------------------------------------------
        # 4. CLOSING SLIDE
        # ---------------------------------------------------------------
        try:
            closing_layout_name = spec_dict.get("closing_layout_name", "End Slide 1")
            closing_layout = get_smart_layout(prs, closing_layout_name)
            closing = prs.slides.add_slide(closing_layout)
            actual_closing_name = closing_layout.name

            # If the closing layout has no text slots (End Slide 1/2/3), the
            # render call is a no-op. Otherwise we set the closing title.
            closing_slot_map = get_slot_map(actual_closing_name)
            if not closing_slot_map.is_end_slide and not closing_slot_map.is_blank:
                closing_title = spec_dict.get("closing_title", "Thank You")

                class _ClosingSpec:
                    pass

                closing_spec = _ClosingSpec()
                closing_spec.title = closing_title
                closing_spec.subhead = spec_dict.get("closing_subhead")
                closing_spec.chapter_label = None
                closing_spec.bullets = None

                render_slide_content(
                    closing, closing_spec, actual_closing_name,
                    prs=prs, log=log,
                )
        except Exception as e:
            log.warning(f"Could not generate closing slide: {e}")

        # ---------------------------------------------------------------
        # 5. SAVE
        # ---------------------------------------------------------------
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            return tmp.name
    except Exception as e:
        log.error(f"Render failed: {e}", exc_info=True)
        return f"Error: Render failed. {e}"


# =============================================================================
# generate_and_render_deck — agent-facing tool entry point
# =============================================================================

async def generate_and_render_deck(
    tool_context: ToolContext,
    deck_spec: dict | None = None,
    spec_artifact_name: str | None = None,
    template_path: str | None = None,
) -> dict:
    """
    Orchestrate the entire deck generation process.
    """
    log = get_logger("generate_and_render_deck_tool")
    try:
        spec_dict = deck_spec

        # 1. PRIORITY: Load from Session State
        if not spec_dict and not spec_artifact_name:
            spec_dict = tool_context.state.get("current_deck_spec")
            if spec_dict:
                log.info("Loaded DeckSpec from session state.")

        # 2. FALLBACK: Load from Artifact Store
        if not spec_dict and spec_artifact_name:
            log.info(f"Loading DeckSpec from artifact: '{spec_artifact_name}'")
            try:
                artifact = await tool_context.load_artifact(spec_artifact_name)
                if not artifact:
                    log.warning(
                        f"Artifact '{spec_artifact_name}' not found. Waiting 2s..."
                    )
                    await asyncio.sleep(2.0)
                    artifact = await tool_context.load_artifact(spec_artifact_name)

                if artifact:
                    spec_json = (
                        artifact.inline_data.data
                        if isinstance(artifact, types.Part)
                        else artifact
                    )
                    if isinstance(spec_json, (bytes, bytearray)):
                        spec_dict = json.loads(spec_json.decode("utf-8"))
                    elif isinstance(spec_json, str):
                        spec_dict = json.loads(spec_json)
                    else:
                        spec_dict = spec_json
            except Exception as e:
                log.error(f"Failed to load named spec artifact: {e}")

        if not spec_dict:
            return {
                "status": "Failed",
                "message": (
                    "No active presentation plan found in session state. "
                    "Please provide deck_spec or ensure an outline was generated."
                ),
            }

        # 3. GCS-FALLBACK TEMPLATE RECOVERY
        working_template = template_path
        if not working_template or not os.path.exists(working_template):
            log.info("Template path invalid or lost. Re-downloading from GCS...")
            working_template = await get_gcs_file_as_local_path(
                DEFAULT_TEMPLATE_URI
            )

        # 4. STANDARDIZE STRUCTURE
        if isinstance(spec_dict.get("slides"), dict):
            spec_dict["slides"] = list(spec_dict["slides"].values())
        if "closing_title" not in spec_dict:
            spec_dict["closing_title"] = "Thank You"

        validated_spec = DeckSpec(**spec_dict)

        all_content = [validated_spec.cover] + validated_spec.slides

        # 5. VISUAL BUDGET & LAYOUT COERCION
        # Allow up to 5 visuals. When a slide has a visual, coerce its layout
        # to one of the Capgemini layouts that has a picture placeholder.
        hard_limit = 5
        visuals_kept = 0
        for slide in validated_spec.slides:
            if slide.visual_prompt:
                if visuals_kept < hard_limit:
                    visuals_kept += 1
                    if slide.layout_name not in LAYOUTS_WITH_PICTURE:
                        log.info(
                            f"Coercing slide '{slide.title}' from "
                            f"'{slide.layout_name}' to '{DEFAULT_IMAGE_LAYOUT}' "
                            f"because it has a visual."
                        )
                        slide.layout_name = DEFAULT_IMAGE_LAYOUT
                else:
                    # Past the visual budget — strip the visual.
                    slide.visual_prompt = None

        # 6. GENERATE VISUALS IN PARALLEL
        tasks = []
        slides_with_visuals = []
        for item in all_content:
            if hasattr(item, "visual_prompt") and item.visual_prompt:
                tasks.append(
                    asyncio.create_task(
                        asyncio.wait_for(
                            generate_visual(item.visual_prompt), timeout=60.0
                        )
                    )
                )
                slides_with_visuals.append(item)

        images = await asyncio.gather(*tasks, return_exceptions=True)
        for s, img in zip(slides_with_visuals, images):
            if not isinstance(img, Exception):
                s.image_data = img

        # 7. RENDER & SAVE
        out_name = f"{validated_spec.cover.title}_{uuid.uuid4().hex[:6]}.pptx"

        local_path = await render_deck_from_spec(
            validated_spec.model_dump(),
            out_name,
            tool_context,
            working_template,
        )
        if local_path.startswith("Error:"):
            return {"status": "Failed", "message": local_path}

        msg = await save_presentation(
            tool_context, out_name, local_path, GCS_BUCKET_NAME
        )
        os.remove(local_path)
        return {"status": "Success", "message": msg}
    except Exception as e:
        log.error(f"Generation failed: {e}", exc_info=True)
        return {"status": "Failed", "message": str(e)}
